from __future__ import annotations

import csv
import hashlib
import json
import mimetypes
import re
import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Iterable
from xml.etree import ElementTree


class KnowledgeIngestionError(RuntimeError):
    """Raised when a source cannot be parsed into the normalized knowledge model."""


@dataclass(slots=True)
class ExtractedAsset:
    asset_type: str
    mime_type: str | None = None
    storage_key: str | None = None
    label: str = ""
    page_start: int | None = None
    page_end: int | None = None
    ocr_text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    parent_ordinal: int | None = None


@dataclass(slots=True)
class ExtractedBlock:
    block_type: str
    text: str
    heading_path: str = ""
    page_start: int | None = None
    page_end: int | None = None
    content_type: str = "text"
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(slots=True)
class ExtractionResult:
    blocks: list[ExtractedBlock]
    assets: list[ExtractedAsset] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)


_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".yaml", ".yml", ".toml",
    ".ini", ".cfg", ".conf", ".py", ".js", ".mjs", ".cjs", ".ts", ".tsx",
    ".jsx", ".java", ".kt", ".kts", ".c", ".h", ".cpp", ".hpp", ".cs",
    ".go", ".rs", ".sql", ".sh", ".bat", ".cmd", ".ps1", ".css", ".scss",
}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
_SUPPORTED_EXTENSIONS = _TEXT_EXTENSIONS | {
    ".pdf", ".docx", ".xlsx", ".xlsm", ".csv", ".tsv", ".pptx", ".html",
    ".htm", ".json", ".xml",
} | _IMAGE_EXTENSIONS


def supported_formats() -> list[str]:
    return sorted(_SUPPORTED_EXTENSIONS)


def guess_source_type(filename: str, mime_type: str | None = None) -> str:
    extension = Path(filename or "").suffix.lower()
    if extension:
        return extension.lstrip(".")
    guessed = mime_type or mimetypes.guess_type(filename or "")[0]
    return (guessed or "unknown").split("/", 1)[-1]


def source_sha256(path: Path, *, chunk_size: int = 1024 * 1024) -> str:
    digest = hashlib.sha256()
    with Path(path).open("rb") as handle:
        for chunk in iter(lambda: handle.read(chunk_size), b""):
            digest.update(chunk)
    return digest.hexdigest()


def safe_filename(value: str) -> str:
    name = Path(value or "document").name.strip() or "document"
    name = re.sub(r"[^A-Za-z0-9._()\- ]+", "_", name).strip(" .")
    return (name or "document")[:180]


def _decode_text(raw: bytes) -> str:
    if raw.startswith(b"\xef\xbb\xbf"):
        return raw.decode("utf-8-sig", errors="replace")
    for encoding in ("utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _clean_text(value: Any) -> str:
    text = str(value or "").replace("\x00", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def _table_text(rows: Iterable[Iterable[Any]], *, title: str = "") -> tuple[str, int, int]:
    normalized: list[list[str]] = []
    width = 0
    for row in rows:
        values = [_clean_text(cell) for cell in row]
        if not any(values):
            continue
        normalized.append(values)
        width = max(width, len(values))
    if not normalized:
        return "", 0, 0
    padded = [row + [""] * (width - len(row)) for row in normalized]
    lines: list[str] = []
    if title:
        lines.append(title)
    header = padded[0]
    lines.append(" | ".join(header))
    lines.append(" | ".join(["---"] * width))
    lines.extend(" | ".join(row) for row in padded[1:])
    return "\n".join(lines).strip(), len(padded), width


def _heading_path(levels: list[str]) -> str:
    return " > ".join(item for item in levels if item)


def _markdown_blocks(text: str, *, content_type: str = "text") -> list[ExtractedBlock]:
    lines = text.splitlines()
    headings = [""] * 6
    buffer: list[str] = []
    blocks: list[ExtractedBlock] = []

    def flush() -> None:
        nonlocal buffer
        body = _clean_text("\n".join(buffer))
        if body:
            blocks.append(
                ExtractedBlock(
                    block_type="section",
                    text=body,
                    heading_path=_heading_path(headings),
                    content_type=content_type,
                )
            )
        buffer = []

    for line in lines:
        match = re.match(r"^\s{0,3}(#{1,6})\s+(.+?)\s*$", line)
        if match:
            flush()
            level = len(match.group(1))
            headings[level - 1] = _clean_text(match.group(2))
            for index in range(level, len(headings)):
                headings[index] = ""
            continue
        buffer.append(line)
    flush()
    if not blocks and _clean_text(text):
        blocks.append(ExtractedBlock(block_type="document", text=_clean_text(text), content_type=content_type))
    return blocks


def _parse_plain(path: Path) -> ExtractionResult:
    raw = path.read_bytes()
    text = _decode_text(raw)
    extension = path.suffix.lower()
    content_type = "code" if extension in {
        ".py", ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx", ".java", ".kt", ".kts",
        ".c", ".h", ".cpp", ".hpp", ".cs", ".go", ".rs", ".sql", ".sh", ".bat", ".cmd",
        ".ps1", ".css", ".scss",
    } else "text"
    if extension in {".md", ".markdown"}:
        blocks = _markdown_blocks(text, content_type=content_type)
    else:
        blocks = [ExtractedBlock(block_type="document", text=_clean_text(text), content_type=content_type)]
    return ExtractionResult(blocks=[block for block in blocks if block.text], metadata={"encoding": "auto"})


def _parse_json(path: Path) -> ExtractionResult:
    text = _decode_text(path.read_bytes())
    try:
        data = json.loads(text)
        text = json.dumps(data, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        pass
    return ExtractionResult([ExtractedBlock("document", _clean_text(text), content_type="structured")])


def _parse_xml(path: Path) -> ExtractionResult:
    text = _decode_text(path.read_bytes())
    try:
        root = ElementTree.fromstring(text)
        parts: list[str] = []
        for element in root.iter():
            value = _clean_text(element.text)
            if value:
                parts.append(f"{element.tag}: {value}")
        normalized = "\n".join(parts)
    except ElementTree.ParseError:
        normalized = text
    return ExtractionResult([ExtractedBlock("document", _clean_text(normalized), content_type="structured")])


def _parse_html(path: Path) -> ExtractionResult:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:  # pragma: no cover - dependency is shipped in release requirements
        raise KnowledgeIngestionError("HTML parser dependency beautifulsoup4 is unavailable") from exc
    soup = BeautifulSoup(_decode_text(path.read_bytes()), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    headings = [""] * 6
    blocks: list[ExtractedBlock] = []
    table_number = 0
    for node in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "table"]):
        name = str(node.name).lower()
        if name.startswith("h") and len(name) == 2 and name[1].isdigit():
            level = max(1, min(6, int(name[1])))
            headings[level - 1] = _clean_text(node.get_text(" ", strip=True))
            for index in range(level, 6):
                headings[index] = ""
            continue
        if name == "table":
            rows = [[cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])] for row in node.find_all("tr")]
            table_number += 1
            text, row_count, column_count = _table_text(rows, title=f"Table {table_number}")
            if text:
                blocks.append(ExtractedBlock("table", text, _heading_path(headings), content_type="table", metadata={"rows": row_count, "columns": column_count}))
            continue
        text = _clean_text(node.get_text(" ", strip=True))
        if text:
            blocks.append(ExtractedBlock("section", text, _heading_path(headings), content_type="code" if name == "pre" else "text"))
    return ExtractionResult(blocks, metadata={"title": _clean_text(soup.title.string) if soup.title and soup.title.string else ""})


def _parse_csv(path: Path) -> ExtractionResult:
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    blocks: list[ExtractedBlock] = []
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        header: list[str] | None = None
        group: list[list[str]] = []
        group_start = 2
        row_number = 0
        for row_number, row in enumerate(reader, start=1):
            values = [_clean_text(cell) for cell in row]
            if header is None:
                header = values
                continue
            group.append(values)
            if len(group) >= 80:
                text, rows, columns = _table_text([header, *group], title=f"Rows {group_start}-{row_number}")
                if text:
                    blocks.append(ExtractedBlock("table", text, heading_path=path.stem, content_type="table", metadata={"row_start": group_start, "row_end": row_number, "rows": rows, "columns": columns}))
                group = []
                group_start = row_number + 1
        if header is None:
            return ExtractionResult([])
        if group:
            text, rows, columns = _table_text([header, *group], title=f"Rows {group_start}-{row_number}")
            if text:
                blocks.append(ExtractedBlock("table", text, heading_path=path.stem, content_type="table", metadata={"row_start": group_start, "row_end": row_number, "rows": rows, "columns": columns}))
    return ExtractionResult(blocks, metadata={"delimiter": "tab" if delimiter == "\t" else "comma"})


def _parse_xlsx(path: Path) -> ExtractionResult:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise KnowledgeIngestionError("Excel parser dependency openpyxl is unavailable") from exc
    workbook = load_workbook(path, read_only=True, data_only=False)
    blocks: list[ExtractedBlock] = []
    try:
        for sheet in workbook.worksheets:
            rows_iter = sheet.iter_rows(values_only=True)
            header: list[str] | None = None
            group: list[list[str]] = []
            group_start = 2
            row_number = 0
            for row_number, row in enumerate(rows_iter, start=1):
                values = [_clean_text(cell) for cell in row]
                if not any(values):
                    continue
                if header is None:
                    header = values
                    continue
                group.append(values)
                if len(group) >= 60:
                    text, rows, columns = _table_text([header, *group], title=f"{sheet.title} rows {group_start}-{row_number}")
                    blocks.append(ExtractedBlock("table", text, sheet.title, content_type="table", metadata={"sheet": sheet.title, "row_start": group_start, "row_end": row_number, "rows": rows, "columns": columns}))
                    group = []
                    group_start = row_number + 1
            if header is not None and group:
                text, rows, columns = _table_text([header, *group], title=f"{sheet.title} rows {group_start}-{row_number}")
                blocks.append(ExtractedBlock("table", text, sheet.title, content_type="table", metadata={"sheet": sheet.title, "row_start": group_start, "row_end": row_number, "rows": rows, "columns": columns}))
            elif header is not None and not group and row_number == 1:
                text, rows, columns = _table_text([header], title=sheet.title)
                if text:
                    blocks.append(ExtractedBlock("table", text, sheet.title, content_type="table", metadata={"sheet": sheet.title, "rows": rows, "columns": columns}))
    finally:
        workbook.close()
    return ExtractionResult(blocks, metadata={"sheets": workbook.sheetnames})


def _extract_docx_images(document: Any, asset_dir: Path) -> list[ExtractedAsset]:
    assets: list[ExtractedAsset] = []
    asset_dir.mkdir(parents=True, exist_ok=True)
    seen: set[str] = set()
    counter = 0
    for rel in document.part.rels.values():
        target = getattr(rel, "target_part", None)
        content_type = str(getattr(target, "content_type", "") or "")
        if not content_type.startswith("image/") or target is None:
            continue
        blob = bytes(getattr(target, "blob", b"") or b"")
        digest = hashlib.sha256(blob).hexdigest()
        if not blob or digest in seen:
            continue
        seen.add(digest)
        counter += 1
        extension = mimetypes.guess_extension(content_type) or ".img"
        output = asset_dir / f"image-{counter:04d}{extension}"
        output.write_bytes(blob)
        ocr_text = _ocr_path(output)
        assets.append(ExtractedAsset("image", content_type, storage_key=output.name, label=f"Image {counter}", ocr_text=ocr_text, metadata={"sha256": digest}))
    return assets


def _parse_docx(path: Path, asset_dir: Path) -> ExtractionResult:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise KnowledgeIngestionError("Word parser dependency python-docx is unavailable") from exc
    document = Document(path)
    blocks: list[ExtractedBlock] = []
    headings = [""] * 9
    try:
        content = document.iter_inner_content()
    except AttributeError:  # pragma: no cover - old python-docx fallback
        content = list(document.paragraphs) + list(document.tables)
    for item in content:
        if hasattr(item, "rows"):
            rows = [[cell.text for cell in row.cells] for row in item.rows]
            text, row_count, column_count = _table_text(rows)
            if text:
                blocks.append(ExtractedBlock("table", text, _heading_path(headings), content_type="table", metadata={"rows": row_count, "columns": column_count}))
            continue
        text = _clean_text(getattr(item, "text", ""))
        if not text:
            continue
        style_name = str(getattr(getattr(item, "style", None), "name", "") or "")
        match = re.match(r"Heading\s+(\d+)", style_name, flags=re.I)
        if match:
            level = max(1, min(len(headings), int(match.group(1))))
            headings[level - 1] = text
            for index in range(level, len(headings)):
                headings[index] = ""
            continue
        blocks.append(ExtractedBlock("section", text, _heading_path(headings), content_type="text", metadata={"style": style_name}))
    assets = _extract_docx_images(document, asset_dir)
    for asset in assets:
        if asset.ocr_text:
            blocks.append(ExtractedBlock("image_ocr", asset.ocr_text, content_type="image_ocr", metadata={"asset": asset.storage_key, "label": asset.label}))
    props = document.core_properties
    return ExtractionResult(blocks, assets, metadata={"author": _clean_text(props.author), "subject": _clean_text(props.subject), "title": _clean_text(props.title)})


def _parse_pptx(path: Path, asset_dir: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:  # pragma: no cover
        raise KnowledgeIngestionError("PowerPoint parser dependency python-pptx is unavailable") from exc
    presentation = Presentation(path)
    blocks: list[ExtractedBlock] = []
    assets: list[ExtractedAsset] = []
    asset_dir.mkdir(parents=True, exist_ok=True)
    image_counter = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_title = _clean_text(slide.shapes.title.text) if slide.shapes.title is not None else f"Slide {slide_number}"
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                text, row_count, column_count = _table_text(rows, title=slide_title)
                if text:
                    blocks.append(ExtractedBlock("table", text, slide_title, slide_number, slide_number, "table", {"slide": slide_number, "rows": row_count, "columns": column_count}))
                continue
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image_counter += 1
                image = shape.image
                extension = "." + str(image.ext or "img")
                output = asset_dir / f"slide-{slide_number:04d}-image-{image_counter:04d}{extension}"
                output.write_bytes(image.blob)
                ocr_text = _ocr_path(output)
                asset = ExtractedAsset("image", image.content_type, output.name, f"{slide_title} image {image_counter}", slide_number, slide_number, ocr_text, {"slide": slide_number})
                assets.append(asset)
                if ocr_text:
                    blocks.append(ExtractedBlock("image_ocr", ocr_text, slide_title, slide_number, slide_number, "image_ocr", {"asset": output.name, "slide": slide_number}))
                continue
            if getattr(shape, "has_text_frame", False):
                text = _clean_text(shape.text)
                if text and text != slide_title:
                    blocks.append(ExtractedBlock("section", text, slide_title, slide_number, slide_number, "text", {"slide": slide_number}))
        if not any(block.page_start == slide_number for block in blocks):
            blocks.append(ExtractedBlock("section", slide_title, slide_title, slide_number, slide_number, "text", {"slide": slide_number}))
    return ExtractionResult(blocks, assets, metadata={"slides": len(presentation.slides)})


def _rapidocr(image: Any) -> str:
    try:
        import numpy as np
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        return ""
    engine = getattr(_rapidocr, "_engine", None)
    if engine is None:
        engine = RapidOCR()
        setattr(_rapidocr, "_engine", engine)
    result, _elapsed = engine(np.asarray(image.convert("RGB")))
    if not result:
        return ""
    return _clean_text("\n".join(str(item[1]) for item in result if len(item) > 1 and item[1]))


def _tesseract_ocr(image: Any) -> str:
    try:
        import pytesseract
    except ImportError:
        return ""
    if not shutil.which("tesseract"):
        return ""
    try:
        return _clean_text(pytesseract.image_to_string(image, config="--psm 6"))
    except Exception:
        return ""


def _ocr_image(image: Any) -> str:
    return _rapidocr(image) or _tesseract_ocr(image)


def _ocr_path(path: Path) -> str:
    try:
        from PIL import Image
    except ImportError:
        return ""
    try:
        with Image.open(path) as image:
            return _ocr_image(image)
    except Exception:
        return ""


def ocr_available() -> bool:
    try:
        import rapidocr_onnxruntime  # noqa: F401
        return True
    except ImportError:
        try:
            import pytesseract  # noqa: F401
        except ImportError:
            return False
        return bool(shutil.which("tesseract"))


def _render_pdf_page(path: Path, page_index: int) -> Any:
    try:
        import pypdfium2 as pdfium
    except ImportError as exc:  # pragma: no cover
        raise KnowledgeIngestionError("Scanned-PDF renderer pypdfium2 is unavailable") from exc
    document = pdfium.PdfDocument(str(path))
    try:
        page = document[page_index]
        bitmap = page.render(scale=2.0)
        return bitmap.to_pil()
    finally:
        document.close()


def _parse_pdf(path: Path) -> ExtractionResult:
    try:
        import pdfplumber
    except ImportError as exc:  # pragma: no cover
        raise KnowledgeIngestionError("PDF parser dependency pdfplumber is unavailable") from exc
    blocks: list[ExtractedBlock] = []
    assets: list[ExtractedAsset] = []
    ocr_pages = 0
    table_count = 0
    with pdfplumber.open(path) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            text = _clean_text(page.extract_text() or "")
            used_ocr = False
            if len(text) < 24 and ocr_available():
                try:
                    image = _render_pdf_page(path, page_number - 1)
                    text = _ocr_image(image)
                    used_ocr = bool(text)
                    if used_ocr:
                        ocr_pages += 1
                        assets.append(ExtractedAsset("page_image", "image/png", None, f"Page {page_number}", page_number, page_number, text, {"render_from_source": True, "page": page_number}))
                except Exception:
                    pass
            if text:
                blocks.append(ExtractedBlock("page", text, f"Page {page_number}", page_number, page_number, "image_ocr" if used_ocr else "text", {"page": page_number, "ocr": used_ocr}))
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            for rows in tables:
                table_count += 1
                table_text, row_count, column_count = _table_text(rows, title=f"Table {table_count}, page {page_number}")
                if table_text:
                    blocks.append(ExtractedBlock("table", table_text, f"Page {page_number}", page_number, page_number, "table", {"page": page_number, "table": table_count, "rows": row_count, "columns": column_count}))
        page_count = len(pdf.pages)
    return ExtractionResult(blocks, assets, metadata={"pages": page_count, "ocr_pages": ocr_pages, "tables": table_count})


def _parse_image(path: Path) -> ExtractionResult:
    try:
        from PIL import Image
    except ImportError as exc:  # pragma: no cover
        raise KnowledgeIngestionError("Image parser dependency Pillow is unavailable") from exc
    with Image.open(path) as image:
        width, height = image.size
        ocr_text = _ocr_image(image)
    asset = ExtractedAsset("image", mimetypes.guess_type(path.name)[0], None, path.name, ocr_text=ocr_text, metadata={"width": width, "height": height, "source_image": True})
    blocks = [ExtractedBlock("image_ocr", ocr_text, path.stem, content_type="image_ocr", metadata={"source_image": True})] if ocr_text else []
    return ExtractionResult(blocks, [asset], metadata={"width": width, "height": height, "ocr": bool(ocr_text)})


def parse_document(path: Path, asset_dir: Path) -> ExtractionResult:
    path = Path(path)
    extension = path.suffix.lower()
    if extension not in _SUPPORTED_EXTENSIONS:
        raise KnowledgeIngestionError(f"Unsupported knowledge file type: {extension or 'unknown'}")
    if extension in _TEXT_EXTENSIONS:
        return _parse_plain(path)
    if extension == ".json":
        return _parse_json(path)
    if extension == ".xml":
        return _parse_xml(path)
    if extension in {".html", ".htm"}:
        return _parse_html(path)
    if extension in {".csv", ".tsv"}:
        return _parse_csv(path)
    if extension in {".xlsx", ".xlsm"}:
        return _parse_xlsx(path)
    if extension == ".docx":
        return _parse_docx(path, asset_dir)
    if extension == ".pptx":
        return _parse_pptx(path, asset_dir)
    if extension == ".pdf":
        return _parse_pdf(path)
    if extension in _IMAGE_EXTENSIONS:
        return _parse_image(path)
    raise KnowledgeIngestionError(f"No parser registered for {extension}")


def _split_long_segment(text: str, max_chars: int, overlap_chars: int) -> list[str]:
    text = _clean_text(text)
    if len(text) <= max_chars:
        return [text] if text else []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chars)
        if end < len(text):
            boundary = max(text.rfind("\n", start + max_chars // 2, end), text.rfind(". ", start + max_chars // 2, end))
            if boundary > start:
                end = boundary + 1
        chunk = _clean_text(text[start:end])
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = max(start + 1, end - overlap_chars)
    return chunks


def chunk_block(block: ExtractedBlock, *, max_chars: int = 1500, overlap_chars: int = 180) -> list[dict[str, Any]]:
    prefix = f"{block.heading_path}\n\n" if block.heading_path else ""
    text = _clean_text(block.text)
    if not text:
        return []
    pieces = _split_long_segment(text, max_chars=max_chars, overlap_chars=overlap_chars)
    chunks: list[dict[str, Any]] = []
    for piece in pieces:
        searchable = _clean_text(prefix + piece)
        # Token count is deliberately an estimate in Phase 2. Phase 3 can replace
        # it with the embedding/tokenizer-specific count without changing storage.
        token_count = max(1, (len(searchable) + 3) // 4)
        chunks.append({
            "content_type": block.content_type,
            "text": searchable,
            "token_count": token_count,
            "page_start": block.page_start,
            "page_end": block.page_end,
            "metadata": {**block.metadata, "heading_path": block.heading_path},
            "lexical_status": "pending",
            "vector_status": "pending",
        })
    return chunks


def normalize_result(result: ExtractionResult) -> list[dict[str, Any]]:
    normalized: list[dict[str, Any]] = []
    for ordinal, block in enumerate(result.blocks):
        text = _clean_text(block.text)
        if not text:
            continue
        normalized.append({
            "block_type": block.block_type,
            "ordinal": ordinal,
            "heading_path": block.heading_path,
            "page_start": block.page_start,
            "page_end": block.page_end,
            "text": text,
            "metadata": block.metadata,
            "chunks": chunk_block(block),
        })
    return normalized
