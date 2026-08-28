from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.config import Settings
from app.db import Database
from app.knowledge import KnowledgeEngine
from app.migrations import CURRENT_SCHEMA_VERSION


def _build(tmp_path: Path) -> tuple[Database, KnowledgeEngine, dict]:
    settings = Settings(
        db_path=tmp_path / "verbanode.db",
        backup_path=tmp_path / "backups",
        knowledge_path=tmp_path / "knowledge",
        open_browser=False,
    )
    db = Database(settings)
    db.initialize()
    engine = KnowledgeEngine(db, settings.knowledge_dir)
    library = engine.create_library({"name": "Phase 2", "description": "", "enabled": True})
    return db, engine, library


def _stage_and_ingest(engine: KnowledgeEngine, library_id: int, source: Path):
    staged = engine.new_upload_path(source.name)
    staged.write_bytes(source.read_bytes())
    document, job = engine.register_staged_upload(
        library_id=library_id,
        staged_path=staged,
        source_name=source.name,
        mime_type=None,
    )
    result = engine.ingest_document(int(document["id"]), int(job["id"]))
    return result, engine.document_content(int(document["id"]))


def test_phase2_schema_and_status(tmp_path: Path) -> None:
    db, engine, _library = _build(tmp_path)
    assert CURRENT_SCHEMA_VERSION == 13
    assert db.schema_version() == 13
    status = engine.status()
    assert status["phase"] == "chat_voice_cutover"
    assert status["ingestion_enabled"] is True
    assert status["retrieval_enabled"] is True
    assert status["capabilities"]["tables"] is True
    assert status["capabilities"]["vlm"] is False
    assert ".pdf" in status["supported_formats"]
    assert ".docx" in status["supported_formats"]
    assert ".xlsx" in status["supported_formats"]
    assert ".pptx" in status["supported_formats"]


def test_plain_text_ingestion_builds_parent_child_content(tmp_path: Path) -> None:
    _db, engine, library = _build(tmp_path)
    source = tmp_path / "manual.md"
    source.write_text(
        "# Motor Manual\n\n## Reset Procedure\n\nTurn the controller off before reset.\n\n"
        "Error VN-AE-104 is related to capture initialization.\n",
        encoding="utf-8",
    )
    document, content = _stage_and_ingest(engine, library["id"], source)
    assert document["status"] == "parsed"
    assert content["parent_blocks"]
    assert content["chunks"]
    assert any("Reset Procedure" in chunk["text"] for chunk in content["chunks"])
    assert all(chunk["lexical_status"] == "ready" for chunk in content["chunks"])
    assert all(chunk["vector_status"] in {"ready", "error"} for chunk in content["chunks"])


def test_docx_tables_and_headings_are_preserved(tmp_path: Path) -> None:
    _db, engine, library = _build(tmp_path)
    source = tmp_path / "service.docx"
    doc = Document()
    doc.add_heading("Electrical", level=1)
    doc.add_paragraph("Motor controller specifications.")
    table = doc.add_table(rows=3, cols=3)
    values = [
        ("Motor", "Voltage", "Current"),
        ("Left drive", "24 V", "8 A"),
        ("Head servo", "12 V", "2 A"),
    ]
    for row, data in zip(table.rows, values):
        for cell, value in zip(row.cells, data):
            cell.text = value
    doc.save(source)
    document, content = _stage_and_ingest(engine, library["id"], source)
    assert document["status"] == "parsed"
    assert any(block["heading_path"] == "Electrical" for block in content["parent_blocks"])
    assert any(block["block_type"] == "table" and "Head servo" in block["text"] for block in content["parent_blocks"])


def test_xlsx_and_pptx_ingestion_keep_structure(tmp_path: Path) -> None:
    _db, engine, library = _build(tmp_path)

    workbook_path = tmp_path / "catalog.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Products"
    sheet.append(["SKU", "Name", "Price"])
    sheet.append(["XR4-A", "Shoulder motor", 650])
    sheet.append(["XR4-B", "Head servo", 120])
    workbook.save(workbook_path)
    xlsx_document, xlsx_content = _stage_and_ingest(engine, library["id"], workbook_path)
    assert xlsx_document["status"] == "parsed"
    assert any("Products" in block["heading_path"] and "XR4-A" in block["text"] for block in xlsx_content["parent_blocks"])

    pptx_path = tmp_path / "training.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Safety Procedure"
    slide.placeholders[1].text = "Disconnect 24 V power before servicing the motor controller."
    presentation.save(pptx_path)
    pptx_document, pptx_content = _stage_and_ingest(engine, library["id"], pptx_path)
    assert pptx_document["status"] == "parsed"
    assert any("Safety Procedure" in block["heading_path"] for block in pptx_content["parent_blocks"])
    assert any("Disconnect 24 V" in block["text"] for block in pptx_content["parent_blocks"])


def test_reingest_replaces_content_and_delete_removes_source(tmp_path: Path) -> None:
    _db, engine, library = _build(tmp_path)
    source = tmp_path / "notes.txt"
    source.write_text("first version", encoding="utf-8")
    document, content = _stage_and_ingest(engine, library["id"], source)
    document_id = int(document["id"])
    storage = engine.root / document["storage_key"]
    assert storage.exists()
    storage.write_text("second version", encoding="utf-8")
    job = engine.reingest_document(document_id)
    engine.ingest_document(document_id, int(job["id"]))
    updated = engine.document_content(document_id)
    assert any("second version" in chunk["text"] for chunk in updated["chunks"])
    assert not any("first version" in chunk["text"] for chunk in updated["chunks"])
    engine.delete_document(document_id)
    assert not storage.exists()


def test_pdf_html_and_csv_ingestion(tmp_path: Path) -> None:
    _db, engine, library = _build(tmp_path)

    from reportlab.pdfgen import canvas

    pdf_path = tmp_path / "manual.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 760, "XR4 Motor Troubleshooting")
    pdf.drawString(72, 735, "Error VN-AE-104 indicates capture initialization failure.")
    pdf.save()
    pdf_document, pdf_content = _stage_and_ingest(engine, library["id"], pdf_path)
    assert pdf_document["status"] == "parsed"
    assert any("VN-AE-104" in block["text"] for block in pdf_content["parent_blocks"])
    assert pdf_document["metadata"]["pages"] == 1

    html_path = tmp_path / "guide.html"
    html_path.write_text(
        "<html><head><title>Guide</title></head><body><h1>Electrical</h1>"
        "<p>Disconnect power before service.</p>"
        "<table><tr><th>Motor</th><th>Voltage</th></tr>"
        "<tr><td>Left drive</td><td>24 V</td></tr></table></body></html>",
        encoding="utf-8",
    )
    html_document, html_content = _stage_and_ingest(engine, library["id"], html_path)
    assert html_document["status"] == "parsed"
    assert any(block["block_type"] == "table" and "24 V" in block["text"] for block in html_content["parent_blocks"])

    csv_path = tmp_path / "products.csv"
    csv_path.write_text("SKU,Name,Price\nXR4-A,Shoulder motor,650\nXR4-B,Head servo,120\n", encoding="utf-8")
    csv_document, csv_content = _stage_and_ingest(engine, library["id"], csv_path)
    assert csv_document["status"] == "parsed"
    assert any(block["block_type"] == "table" and "XR4-B" in block["text"] for block in csv_content["parent_blocks"])


def test_dashboard_conversation_rail_has_no_native_scrollbar() -> None:
    root = Path(__file__).resolve().parents[1]
    css = (root / "app" / "static" / "styles.css").read_text(encoding="utf-8")
    assert "v0.10.1 fixed-viewport dashboard polish" in css
    assert ".control-panel {\n  min-height:0;\n  overflow:hidden;" in css
    assert ".control-card {\n  min-height:0;\n  max-height:none;\n  overflow:hidden;" in css
    assert "*::-webkit-scrollbar { width:0; height:0; display:none; }" in css
